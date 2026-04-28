"""Build slide 1 — Hook / Cover."""
import sys, os
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

from _design import (
    SLIDE_W, SLIDE_H, MARGIN_L, CONTENT_W,
    HEADLINE_Y, HEADLINE_H,
    COL_HEAD, COL_SUB, COL_LABEL,
    KOR_FONT, SZ_HEADLINE, SZ_QUOTE, SZ_PAGENUM,
    new_slide, add_text, add_footer,
)

OUT_DIR  = os.path.dirname(__file__)
OUT_PPTX = os.path.join(OUT_DIR, "초안v3.pptx")
X_IMG    = os.path.join(OUT_DIR, "_x_composite.png")
GIST_IMG = r"C:\home\labs\llm-wiki\_excalidraw_images\9f406cb4d08c.png"

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

slide = new_slide(prs)

# ── Headline ─────────────────────────────────────────────────────
add_text(
    slide,
    MARGIN_L, HEADLINE_Y, CONTENT_W, HEADLINE_H,
    "LLM이 일회성으로 매번 답을 찾는 시대에서,\n지식을 스스로 쌓고 진화하는 시대로.",
    size=SZ_HEADLINE, bold=True, color=COL_HEAD,
)

# ── Karpathy name ────────────────────────────────────────────────
KARP_NAME_Y = HEADLINE_Y + HEADLINE_H + Inches(0.08)
add_text(
    slide,
    MARGIN_L, KARP_NAME_Y, CONTENT_W, Inches(0.35),
    "andrej karpathy  (전 Tesla AI 디렉터, OpenAI 창립멤버)",
    size=13, bold=False, color=COL_SUB,
)

# ── Karpathy quote ───────────────────────────────────────────────
KARP_QUOTE_Y = KARP_NAME_Y + Inches(0.37)
add_text(
    slide,
    MARGIN_L, KARP_QUOTE_Y, CONTENT_W, Inches(0.4),
    "“코드 조작보다 더 많은 토큰을 쓰고 있다. 꽤 잘 작동한다.”",
    size=15, bold=True, color=COL_HEAD,
)

# ── Image labels ─────────────────────────────────────────────────
LBL_Y   = KARP_QUOTE_Y + Inches(0.48)
COL_W   = Inches(5.85)
RIGHT_X = MARGIN_L + COL_W + Inches(0.4)

add_text(slide, MARGIN_L, LBL_Y, COL_W, Inches(0.3),
         "LLM Knowledge Bases",
         size=11, bold=True, color=COL_LABEL)
add_text(slide, RIGHT_X, LBL_Y, COL_W, Inches(0.3),
         "LLM Wiki",
         size=11, bold=True, color=COL_LABEL)

# ── Images ───────────────────────────────────────────────────────
def img_hw(path):
    with Image.open(path) as im:
        return im.size   # (w, h)

xw, xh = img_hw(X_IMG)
gw, gh = img_hw(GIST_IMG)

IMG_Y     = LBL_Y + Inches(0.34)
x_img_h   = Inches(COL_W.inches * xh / xw)
g_img_h   = Inches(COL_W.inches * gh / gw)

slide.shapes.add_picture(X_IMG,    MARGIN_L, IMG_Y, COL_W, x_img_h)
slide.shapes.add_picture(GIST_IMG, RIGHT_X,  IMG_Y, COL_W, g_img_h)

# ── Footer ───────────────────────────────────────────────────────
add_footer(slide, pagenum=1)

prs.save(OUT_PPTX)
print("saved:", OUT_PPTX)
print(f"X img height: {x_img_h.inches:.3f}\"  Gist: {g_img_h.inches:.3f}\"")
