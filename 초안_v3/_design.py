"""Shared design constants and helpers for all slide build scripts."""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Slide dimensions ─────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Layout grid ──────────────────────────────────────────────────
MARGIN_L     = Inches(0.7)
CONTENT_W    = Inches(11.93)   # 13.333 - 2 * 0.7
HEADLINE_Y   = Inches(0.45)
HEADLINE_H   = Inches(1.0)
DIVIDER_Y    = Inches(1.52)    # top divider on content slides
CONTENT_Y    = Inches(1.65)    # first content element below divider
FOOTER_Y     = Inches(6.88)
PAGENUM_Y    = Inches(7.1)

# ── Colors ───────────────────────────────────────────────────────
COL_BG    = RGBColor(0xFA, 0xFA, 0xFA)  # off-white background
COL_HEAD  = RGBColor(0x1E, 0x27, 0x61)  # midnight navy
COL_SUB   = RGBColor(0x47, 0x55, 0x69)  # muted slate
COL_LABEL = RGBColor(0x94, 0xA3, 0xB8)  # light gray
COL_DIV   = RGBColor(0xE2, 0xE8, 0xF0)  # divider line

# ── Typography ───────────────────────────────────────────────────
KOR_FONT      = "맑은 고딕"
SZ_HEADLINE   = 28
SZ_SCENARIO   = 13   # scenario label (bold)
SZ_BODY       = 11   # diagram / flow content
SZ_QUOTE      = 10   # karpathy quote / footer
SZ_PAGENUM    = 10


# ── Helpers ──────────────────────────────────────────────────────

def new_slide(prs):
    """Append blank slide with standard background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = COL_BG
    return slide


def add_text(slide, x, y, w, h, text, *,
             size=11, bold=False, italic=False,
             color=COL_SUB, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = KOR_FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_lines(slide, x, y, w, h, lines, *,
              size=SZ_BODY, bold=False, color=COL_SUB,
              line_spacing=Pt(19)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = KOR_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_divider(slide, y=None):
    """Thin horizontal rule. Defaults to standard DIVIDER_Y."""
    if y is None:
        y = DIVIDER_Y
    r = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_L, y, CONTENT_W, Inches(0.025)
    )
    r.line.fill.background()
    r.fill.solid()
    r.fill.fore_color.rgb = COL_DIV


def add_headline(slide, text):
    """Standard headline — same position/size on every slide."""
    add_text(
        slide,
        MARGIN_L, HEADLINE_Y, CONTENT_W, HEADLINE_H,
        text,
        size=SZ_HEADLINE, bold=True, color=COL_HEAD,
    )


def add_footer(slide, pagenum, left_text=None, right_text=None):
    """Page number + optional left/right footer text."""
    add_text(
        slide,
        MARGIN_L, PAGENUM_Y, Inches(4), Inches(0.3),
        f"{pagenum} / 5",
        size=SZ_PAGENUM, color=COL_LABEL,
    )
    if left_text:
        add_text(
            slide,
            MARGIN_L, FOOTER_Y, Inches(5), Inches(0.35),
            left_text,
            size=SZ_QUOTE, bold=True, color=COL_HEAD,
        )
    if right_text:
        add_text(
            slide,
            Inches(4.0), FOOTER_Y, Inches(8.63), Inches(0.35),
            right_text,
            size=SZ_QUOTE, color=COL_LABEL, align=PP_ALIGN.RIGHT,
        )
